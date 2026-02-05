#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#Libraries
import streamlit as st 
import whisper         
import os              
import google.generativeai as GenAI
from pptx import Presentation 
from io import BytesIO
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
import re
#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
# This is the visual part of the page 
st.set_page_config(page_title="Gen", page_icon="🪄")
st.markdown("""
<h1 style="
    color: #FFFFFF;
    text-align: center;
    text-shadow: 2px 2px 10px rgba(0,0,0,0.7);
">
🪄 Transcription and Slide Creator
</h1>
""", unsafe_allow_html=True)
#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
Url_Imagen = "https://i.pinimg.com/originals/cf/a2/39/cfa239195d194b724a9d38362859a1af.jpg"
#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
st.markdown(
    f"""
    <style>
    .stApp {{
        background-image: url("{Url_Imagen}");
        background-size: cover;
        background-position: center;
        background-repeat: no-repeat;
        background-attachment: fixed;
    }}

    /* Capa oscura para que el texto se lea bien */
    .main {{
        background-color: rgba(0, 0, 0, 0.45);
        padding: 20px;
        border-radius: 20px;
    }}
    </style>
    """,
    unsafe_allow_html=True
)
#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#Enter your API KEY here 
GenAI.configure(api_key=st.secrets["API_KEY"])
#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
# PowerPoint Function
def crear_pptx(texto_generado):
    prs = Presentation()

    # Busca bloques entre --- SLIDE N ---
    pattern = r"---\s*SLIDE\s*\d+\s*---\s*(.*?)\s*(?=(?:---\s*SLIDE\s*\d+\s*---)|\Z)"
    slides = re.findall(pattern, texto_generado, flags=re.S)

    # Fallback si no encuentra el patrón (mantén compatibilidad)
    if not slides:
        slides = [s for s in re.split(r"---\s*SLIDE", texto_generado) if s.strip()]

    for slide_text in slides:
        # Limpia líneas vacías
        lines = [l.strip() for l in slide_text.strip().splitlines() if l.strip()]
        if not lines:
            continue

        # Primera línea = título
        title = lines[0]

        # Buscamos si hay una sección de notas (ej: empieza con "notes" o "notes_slide")
        notes_idx = None
        for i, ln in enumerate(lines[1:], start=1):
            if ln.lower().startswith("notes") or ln.lower().startswith("notes_slide") or ln.lower().startswith("notes:"):
                notes_idx = i
                break

        if notes_idx is not None:
            bullets_lines = lines[1:notes_idx]
            notes_lines = lines[notes_idx+1:]  # lo que venga después de la etiqueta notes
        else:
            bullets_lines = lines[1:]
            notes_lines = []

        # Normaliza bullets: quita prefijos tipo "*", "-", "•"
        bullets = [re.sub(r'^[\*\-\u2022]\s*', '', b) for b in bullets_lines]

        # Crea la diapositiva con layout "Título y contenido"
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Título
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Cuerpo (placeholders[1]) -> formatear como bullets nativos
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()  # limpia contenido por defecto
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            if bullets:
                # primer párrafo
                p = tf.paragraphs[0]
                p.text = bullets[0]
                p.level = 0
                # ajustar tamaño fuente
                for run in p.runs:
                    run.font.size = Pt(20)

                # resto de bullets
                for b in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(18)
            else:
                tf.text = ""


        # Notas del orador (speaker notes) — aquí vamos a escribirlas
        notes_text = ""
        if notes_lines:
            notes_text = "\n".join(notes_lines)
        else:
            # Si no hay sección explícita de notas, usamos los bullets como guía
            notes_text = "\n".join(bullets)

        # Asigna las notas reales
        try:
            slide.notes_slide.notes_text_frame.text = notes_text
        except Exception:
            # en caso raro que no exista, lo ignoramos
            pass

    # Guarda a bytes
    pptx_io = BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()

#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#Transcription Function
Audio_fill = st.file_uploader("Upload your audio so we can transcribe", type=["mp3", "mp4" , "opus" ,"wav", "m4a"])

if Audio_fill is not None:
    st.subheader("🎧Preview your audio")
    st.audio(Audio_fill)


if Audio_fill is not None:

    MAX_FILE_SIZE = 10 * 1024 * 1024
   
    if Audio_fill.size > MAX_FILE_SIZE:
        st.error("The audio is too long or too short. Please upload a file shorter than 3 minutes. (MAX 10MB)")
        st.stop()
    else:
       
        with open("temp_audio.mp3", "wb") as f:
            f.write(Audio_fill.getbuffer())
            
        # We show the loading message so the user can wait.
        with st.spinner("Whisper is processing your audio"):
            modelo_whisper = whisper.load_model("base")
            resultado = modelo_whisper.transcribe("temp_audio.mp3")

    st.success("Transcription success")
    with st.expander("Show transcription"):
        st.write(resultado["text"])

    if st.button("✨ Generative Slides"):
        
        with st.spinner("Gemini is creating your slides..."):
          
            modelo_gemini = GenAI.GenerativeModel('models/gemini-2.5-flash')
            
            instruction = f"""
  
            Analyze the audio transcript: {resultado['text']} and generate ONLY clearly separated slides following these STRICT rules.

            !!! CRITICAL: LANGUAGE ENFORCEMENT !!!
            1. FIRST, analyze the input text to identify the source language exactly.
            2. YOUR OUTPUT MUST BE 100% IN THAT IDENTIFIED SOURCE LANGUAGE.
            3. IF the audio is in English -> Generate slides/notes in ENGLISH.
            4. IF the audio is in French -> Generate slides/notes in FRENCH.
            5. If it's another language, do the same.
            6. DO NOT translate to Spanish unless the audio is actually in Spanish.
            
            === BEGIN DESIGN & CONTENT INSTRUCTIONS ===
            
            1. GOAL
            Create a visually engaging, well-structured presentation based on the audio.
            Avoid walls of text. Use "Visual Markdown" to make it look professional.
            
            2. INSTRUCTION DETECTION
            Determine whether the audio contains a clear instruction to create content.
            
            3. IF A CLEAR INSTRUCTION EXISTS
            Generate a presentation with a MINIMUM of 5 SLIDES.
            !!! CRITICAL SEPARATION RULE !!!
            You MUST insert a horizontal rule (---) between every single slide. 
            Do not write all content in a continuous block. Each slide must be a distinct unit separated by "---".
            
            4. SLIDE STRUCTURE (MANDATORY & VISUAL)
            Each slide MUST follow this exact internal structure:
            
            [INSERT RELEVANT EMOJI] TITLE OF THE SLIDE
            
            Visual Concept: [Describe in 1 sentence a suggestion for an image or icon that fits this slide]
            
            🔹 **[Keyword or Main Idea]:** [Explanation text]
            🔸 **[Keyword or Main Idea]:** [Explanation text]
            🔹 **[Keyword or Main Idea]:** [Explanation text]
            
            notes_slide:
            [Full, natural speaker notes written as if a real presenter were explaining the slide aloud. MUST BE IN THE SAME LANGUAGE AS TRANSCRIPT]
            
            ---
            
            5. FORMATTING RULES FOR "PRETTIER" SLIDES
            - Use Emojis (🔹, 🔸, 🚀, 💡, ✅) as bullet points instead of simple dots.
            - ALWAYS bold the key concept at the start of a bullet point (e.g., **Efficiency:**).
            - Keep bullet points concise (maximum 4 lines per point).
            
            6. IF NO CLEAR INSTRUCTION EXISTS
            Generate ONLY ONE slide stating that an explicit instruction is required (in the source language).
            That slide MUST also include notes_slide.
            
            7. OUTPUT RESTRICTIONS
            - Speaker notes must appear ONLY inside notes_slide.
            - Do NOT place notes in the slide body.
            - ENSURE the "---" separator appears after the notes of every slide except the last one.
            """
#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
            answer = modelo_gemini.generate_content(instruction)
            
            st.header("📝 Generated Content")
            
            st.info("Everything is ready! You can review the content below and download your slides.")
            st.write(answer.text)
            
            pptx_data = crear_pptx(answer.text)
            
        
            st.write("") 
            
            st.download_button(
                label="🚀 DOWNLOAD YOUR POWERPOINT",
                data=pptx_data,
                file_name="Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True 
            )
            st.balloons()

            if os.path.exists("temp_audio.mp3"):
                os.remove("temp_audio.mp3")









